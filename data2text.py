from pptx.dml.color import RGBColor
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.enum.lang import MSO_LANGUAGE_ID
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

from base import ObjectWorker
from utils.pptx_params import textbox

import pandas as pd


class TextWorker(ObjectWorker):
    def __init__(self, prs=False):
        super(ObjectWorker, self).__init__()
        if prs:
            # initializing a basic presentation file when no existing file is given
            self.prs = self.create_prs()

        self.uid_pool = []
        self.alignment = {"left": PP_PARAGRAPH_ALIGNMENT.LEFT,
                          "right": PP_PARAGRAPH_ALIGNMENT.RIGHT,
                          "center": PP_PARAGRAPH_ALIGNMENT.CENTER}
        self.text_lang = {"tc": MSO_LANGUAGE_ID.TRADITIONAL_CHINESE}

    def creat_text(self, uid, data, slide, obj_format, position):
        # create text box
        shapes = slide.shapes
        x, y, w, h = position
        shape = shapes.add_textbox(x, y, w, h)

        text_frame = shape.text_frame
        text_frame.clear()
        p = text_frame.paragraphs[0]

        # text box content is determined here
        p.text = data

        # set the text format
        text_format = self.default_text_format(obj_format)
        textbox_format = self.default_textbox_format(obj_format)
        self.text_format_setter(p, text_format, shape, textbox_format)
        self.uid_pool.append(uid)
        return slide

    def text_format_setter(self, paragraph, text_format, shape, textbox_format):
        # TODO: Replace multiple ifs with mapping table
        # here is for text format
        if 'alignment' in text_format:
            paragraph.alignment = self.alignment[text_format["alignment"]]

        if "font" in text_format:
            paragraph.font.name = textbox(text_format['font'])

        if 'font_size' in text_format:
            paragraph.font.size = Pt(text_format['font_size'])

        if 'font_color' in text_format:
            paragraph.font.color.rgb = text_format['font_color']

        if 'language' in text_format:
            paragraph.font.language_id = self.text_lang[text_format['language']]

        # below is for text box format
        if 'color' in textbox_format:
            if textbox_format['color'] == "no_fill":
                shape.fill.background()

        return paragraph

    @ staticmethod
    def default_text_format(text_format):
        default_format = {"alignment": "left",
                          "font": "title_font",
                          "font_size": 12,
                          'font_color': RGBColor(0, 0, 0)}

        if text_format is None:
            return default_format

        for t in text_format:
            default_format[t] = text_format[t]

        return default_format

    @staticmethod
    def default_textbox_format(text_format):
        # this is for the format of textbox (e.g. color,...)
        default_format = {'color': "no_fill"}

        if text_format is None:
            return default_format

        for t in text_format:
            default_format[t] = text_format[t]

        return default_format

