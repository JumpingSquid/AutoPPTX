"""
Data2chart ver 1.1

Description:
Data2Chart creates chart(s) on presentation object(prs).
prs can be taken as input or o/w initialized at the beginning.
ver 1.1 allows for direct output of prs file with chart, while it's still available
  being called by pptx_construct
For simplicity, Data2chart should not allow for complicated layout manipulation.
To create slide with multiple charts and auto layout management, one should use pptx_construct instead.
"""

from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.chart import XL_LEGEND_POSITION, XL_DATA_LABEL_POSITION
from pptx.util import Pt, Inches
from pptx import Presentation

from utils.pptx_params import PrsParamsManager
from base import ObjectWorker
import pandas


class ChartWorker(ObjectWorker):

    def __init__(self, prs=False):
        super(ObjectWorker, self).__init__()
        if prs:
            # initializing a basic presentation file when no existing file is given
            self.prs = self.create_prs()

            self.slide_pool = {}
            bullet_slide_layout = self.prs.slide_layouts[6]
            self.slide_pool[0] = self.prs.slides.add_slide(bullet_slide_layout)

        # default origin and size for the chart
        self.origin = (Inches(0), Inches(0))
        self.size = (Inches(13), Inches(6))
        self.uid_pool = []


    def add_chart(self, data, slide_id, chart_type, position=None):

        if chart_type not in ["hist", "stackbar", "bar", "pie", "line", "stackcolumn"]:
            print(f"Unknow chart type: {chart_type} is given, will create bar chart directly")
            chart_type = 'bar'

        if isinstance(data, pandas.DataFrame):
            chart_data = self.pandas_to_ppt_chart_data(data)
        else:
            raise TypeError("data type is not supported")

        self.create_chart(chart_data, self.slide_pool[slide_id], chart_type, position)

    def add_slide(self, slide_id, slide_type, overwrite=False):
        if slide_id in self.slide_pool:
            if not overwrite:
                raise ValueError(f"slide id {slide_id} existing")
        bullet_slide_layout = self.prs.slide_layouts[slide_type]
        self.slide_pool[slide_id] = self.prs.slides.add_slide(bullet_slide_layout)

    def save(self, filepath):
        try:
            self.prs.save(filepath)
            return 1
        except Exception:
            return -1

    def create_chart(self, data, slide, obj_format, uid, position=None):
        # called by the add_chart when data2chart is used as an independent module
        # or called by the PrsLayoutManager

        # set the position of the chart
        if position is None:
            x, y = self.origin[0], self.origin[1]
            w, h = self.size[0], self.size[1]
        else:
            x, y, w, h = position

        # analyze the format argument
        chart_format_setter = ChartFormatSetter(obj_format)
        chart_type = chart_format_setter.chart_format['chart_type']

        # add chart to slide
        chart_type_dict = {"hist": XL_CHART_TYPE.COLUMN_CLUSTERED,
                           "stackbar": XL_CHART_TYPE.BAR_STACKED_100,
                           "stackcolumn": XL_CHART_TYPE.COLUMN_STACKED_100,
                           "bar": XL_CHART_TYPE.BAR_CLUSTERED,
                           "pie": XL_CHART_TYPE.PIE,
                           "line": XL_CHART_TYPE.LINE}

        graphic_frame = slide.shapes.add_chart(
            chart_type_dict[chart_type], x, y, w, h, data
        )

        # chart format setting
        if chart_type == "pie":
            chart_format_setter.pie_chart_format(graphic_frame.chart)
        elif chart_type == "line":
            chart_format_setter.line_chart_format(graphic_frame.chart)
        else:
            chart_format_setter.general_chart_format(graphic_frame.chart)

        # store the uid if chart is created successfully
        self.uid_pool.append(uid)

        return slide

    @staticmethod
    # TODO: handled by data_preprocessor in the future
    def pandas_to_ppt_chart_data(dataframe):
        assert isinstance(dataframe, pandas.DataFrame)

        data = CategoryChartData()
        col_lst = dataframe.columns.to_list()
        data.categories = col_lst

        for i, r in dataframe.iterrows():
            data.add_series(i, tuple([r[col] for col in col_lst]))

        return data


class ChartFormatSetter:
    def __init__(self, chart_format=None):
        self.params = PrsParamsManager()
        self.chart_format = self.chart_format_inference(chart_format)

    def general_chart_format(self, chart):
        chart.value_axis.visible = False
        chart.value_axis.has_major_gridlines = False
        chart.font.size = self.chart_format["chart_font_size"]

        # set title
        if self.chart_format["chart_title"] is not None:
            chart.has_title = True
            chart.chart_title.has_text_frame = True
            chart.chart_title.text_frame.text = self.chart_format["chart_title"]
            chart.chart_title.text_frame.paragraphs[0].font.size = Pt(13)

        # set legends
        chart.has_legend = self.chart_format["legend_bool"]
        if self.chart_format["legend_bool"]:
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.font.size = self.chart_format["legend_font_size"]

            chart.legend.include_in_layout = False

        # set data labels
        chart.plots[0].has_data_labels = self.chart_format["label_bool"]
        if self.chart_format["label_bool"]:
            labels = chart.plots[0].data_labels
            labels.number_format_is_linked = False
            labels.number_format = self.chart_format["label_number_format"]
            labels.font.size = self.chart_format["label_font_size"]

        # color
        if self.chart_format["colormap"] is None:
            color_r = 0
            color_g = 0
            color_b = 0
            colorgradient = len(chart.series)
            for series in chart.series:
                fill = series.format.fill  # fill the legend as well
                fill.solid()
                fill.fore_color.rgb = RGBColor(color_r, color_g, color_b)
                color_r += int(255/colorgradient)
                color_g += int(255/colorgradient)
                color_b += int(255/colorgradient)
        else:
            # color map format
            # [[color_R1, color_G1, color_B1], [color_R2, color_G2, color_B2],...]
            colormap = self.chart_format["colormap"]
            colormap_index = 0
            for series in chart.series:
                fill = series.format.fill  # fill the legend as well
                fill.solid()
                fill.fore_color.rgb = RGBColor(colormap[colormap_index % len(colormap)][0],
                                               colormap[colormap_index % len(colormap)][1],
                                               colormap[colormap_index % len(colormap)][2])
                colormap_index += 1
        return

    def line_chart_format(self, chart):
        chart.value_axis.visible = False
        chart.value_axis.has_major_gridlines = False
        chart.font.size = self.chart_format["chart_font_size"]

        # set legends
        chart.has_legend = self.chart_format["legend_bool"]
        if self.chart_format["legend_bool"]:
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.font.size = self.chart_format["legend_font_size"]
            chart.legend.include_in_layout = False

        # set data labels
        chart.plots[0].has_data_labels = self.chart_format["label_bool"]
        if self.chart_format["label_bool"]:
            labels = chart.plots[0].data_labels
            labels.number_format_is_linked = False
            labels.number_format = self.chart_format["label_number_format"]
            labels.font.size = self.chart_format["label_font_size"]

        # color
        if self.chart_format["colormap"] is None:
            color_r = 0
            color_g = 0
            color_b = 0
            colorgradient = len(chart.series)
            for series in chart.series:
                fill = series.format.fill  # fill the legend as well
                fill.solid()
                fill.fore_color.rgb = RGBColor(color_r, color_g, color_b)
                color_r += int(255/colorgradient)
                color_g += int(255/colorgradient)
                color_b += int(255/colorgradient)
        else:
            # color map format
            # [[color_R1, color_G1, color_B1], [color_R2, color_G2, color_B2],...]
            colormap = self.chart_format["colormap"]
            colormap_index = 0
            for series in chart.series:
                fill = series.format.fill  # fill the legend as well
                fill.solid()
                fill.fore_color.rgb = RGBColor(colormap[colormap_index % len(colormap)][0],
                                               colormap[colormap_index % len(colormap)][1],
                                               colormap[colormap_index % len(colormap)][2])
                colormap_index += 1
        return

    def pie_chart_format(self, chart):
        chart.has_legend = self.chart_format["legend_bool"]
        chart.font.size = self.chart_format["chart_font_size"]
        # set title
        if self.chart_format["chart_title"] is not None:
            chart.has_title = True
            chart.chart_title.has_text_frame = True
            chart.chart_title.text_frame.text = self.chart_format["chart_title"]
            chart.chart_title.text_frame.paragraphs[0].font.size = Pt(13)

        if self.chart_format["legend_bool"]:
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False
            chart.legend.font.size = self.chart_format["legend_font_size"]

        chart.plots[0].has_data_labels = self.chart_format["label_bool"]
        if self.chart_format["label_bool"]:
            data_labels = chart.plots[0].data_labels
            data_labels.number_format = self.chart_format["label_number_format"]
            data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END

        if self.chart_format["colormap"] is not None:
            colormap = self.chart_format["colormap"]
            colormap_index = 0
            for point in chart.series[0].points:
                fill = point.format.fill  # fill the legend as well
                fill.solid()
                fill.fore_color.rgb = RGBColor(colormap[colormap_index % len(colormap)][0],
                                               colormap[colormap_index % len(colormap)][1],
                                               colormap[colormap_index % len(colormap)][2])
                colormap_index += 1

    def chart_format_inference(self, predefined_chart_format=None):
        chart_format = self.params.get_chart_format()

        if predefined_chart_format is not None:
            for chart_format_key in predefined_chart_format:
                chart_format[chart_format_key] = predefined_chart_format[chart_format_key]

                if chart_format_key == 'font_size':
                    chart_format['legend_font_size'] = predefined_chart_format[chart_format_key]
                    chart_format['label_font_size'] = predefined_chart_format[chart_format_key]
                    chart_format['chart_font_size'] = predefined_chart_format[chart_format_key]

        return chart_format
