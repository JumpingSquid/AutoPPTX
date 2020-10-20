"""
Data2chart ver 1.1

Description:
Data2Chart creates chart(s) on presentation object(prs).
prs can be taken as input or o/w initialized at the beginning.
ver 1.1 allows for direct output of prs file with chart, while it's still available
  being called by pptx_construct
For simplicity, Data2chart should not allow for complicated layout manipulation,
that is, one slide for one chart rule should be follow.
To create slide with multiple charts, one should use pptx_construct instead.
"""

from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.chart import XL_LEGEND_POSITION, XL_DATA_LABEL_POSITION
from pptx.util import Pt
from pptx import Presentation
from pptx.util import Inches


class ChartCreator:

    def __init__(self,
                 prs=None,
                 origin=None,
                 size=None,
                 chart_category=None,
                 chart_series=None,
                 chart_type="line",
                 chart_format=None
                 ):

        if prs is None:
            # initializing a basic presentation file when no existing file is given
            self.prs = self.create_prs()
        else:
            assert isinstance(Presentation, prs)
            self.prs = prs

        self.slide_pool = {}
        bullet_slide_layout = self.prs.slide_layouts[6]
        self.slide_pool[0] = self.prs.slides.add_slide(bullet_slide_layout)

        if chart_series is not None:
            self.chart_series = chart_series
            if chart_category is not None:
                self.chart_category = chart_category

        if origin is not None:
            self.origin = origin
        else:
            print("Warning: No origin is give, use (0,0) as default")
            self.origin = (Inches(0), Inches(0))

        if size is not None:
            self.size = size
        else:
            print("Warning: No origin is give, use (6 inch,6 inch) as default")
            self.size = (Inches(13), Inches(6))

        # create default chart format
        self.chart_format_setter = ChartFormatSetter()
        self.chart_format = self.chart_format_setter.chart_format_inference()

        # if some params are given, adjust to new
        if chart_format is not None:
            for k in chart_format:
                self.chart_format[k] = chart_format[k]

    def chart_create(self, slide, chart_type):
        # define categorical data
        chart_data = CategoryChartData()
        chart_data.categories = [cat for cat in self.chart_category]

        # define value data -
        for series in self.chart_series:
            chart_data.add_series(series[0], series[1])

        # set the position of the chart
        x, y = self.origin[0], self.origin[1]
        cx, cy = self.size[0], self.size[1]

        # add chart to slide
        chart_type_dict = {"hist": XL_CHART_TYPE.COLUMN_CLUSTERED,
                           "stackbar": XL_CHART_TYPE.BAR_STACKED_100,
                           "stackcolumn": XL_CHART_TYPE.COLUMN_STACKED_100,
                           "bar": XL_CHART_TYPE.BAR_CLUSTERED,
                           "pie": XL_CHART_TYPE.PIE,
                           "line": XL_CHART_TYPE.LINE}

        graphic_frame = slide.shapes.add_chart(
            chart_type_dict[chart_type], x, y, cx, cy, chart_data
        )

        # chart default theme setting
        if chart_type == "pie":
            self.chart_format_setter.piechartformat(graphic_frame.chart)
        elif chart_type == "line":
            self.chart_format_setter.linechartformat(graphic_frame.chart)
        else:
            self.chart_format_setter.chartformat(graphic_frame.chart)

        return slide

    def add_chart(self, data, slide_id, chart_type):
        if chart_type in ["hist", "stackbar", "bar", "pie", "line", "stackcolumn"]:
            return self.chart_create(self.slide_pool[slide_id], chart_type)
        else:
            raise ValueError(f"Unknow chart type: {str(chart_type)} is given")

    def add_slide(self, slide_id, slide_type, overwrite=False):
        assert isinstance(slide_id, int)
        if slide_id in self.slide_pool:
            if not overwrite:
                raise ValueError(f"slide id {slide_id} existing")
        bullet_slide_layout = self.prs.slide_layouts[slide_type]
        self.slide_pool[slide_id] = self.prs.slides.add_slide(bullet_slide_layout)

    def pandas_to_ppt_series(self, data):
        import pandas
        if isinstance(data, pandas.DataFrame):
            chart_series = []
            for col in data.columns:
                chart_series.append((str(col), data[col]))

            self.chart_category = data.index.to_list()
            self.chart_series = chart_series

        elif isinstance(data, pandas.Series):
            self.chart_category = data.index.to_list()
            self.chart_series = [("series_1", tuple(data.values.to_list()))]

    @staticmethod
    def create_prs() -> Presentation():
        prs = Presentation()
        # slide size: 16:9
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        return prs


class ChartFormatSetter:
    def __init__(self):
        self.chart_format = None

    def chartformat(self, chart):
        chart.value_axis.visible = False
        chart.value_axis.has_major_gridlines = False
        chart.font.size = self.chart_format["chart_font_size"]

        # set title
        if self.chart_format["chart_title"] is not None:
            chart.has_title = True
            chart.chart_title.has_text_frame = True
            chart.chart_title.text_frame.text = self.chart_format["chart_title"]
            chart.chart_title.text_frame.paragraphs[0].font.size = Pt(13)

        # set legends
        chart.has_legend = self.chart_format["legend_bool"]
        if self.chart_format["legend_bool"]:
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.font.size = self.chart_format["legend_font_size"]

            chart.legend.include_in_layout = False

        # set data labels
        chart.plots[0].has_data_labels = self.chart_format["label_bool"]
        if self.chart_format["label_bool"]:
            labels = chart.plots[0].data_labels
            labels.number_format_is_linked = False
            labels.number_format = self.chart_format["label_number_format"]
            labels.font.size = self.chart_format["label_font_size"]

        # color
        if self.chart_format["colormap"] is None:
            color_r = 0
            color_g = 0
            color_b = 0
            colorgradient = len(chart.series)
            for series in chart.series:
                fill = series.format.fill  # fill the legend as well
                fill.solid()
                fill.fore_color.rgb = RGBColor(color_r, color_g, color_b)
                color_r += int(255/colorgradient)
                color_g += int(255/colorgradient)
                color_b += int(255/colorgradient)
        else:
            # color map format
            # [[color_R1, color_G1, color_B1], [color_R2, color_G2, color_B2],...]
            colormap = self.chart_format["colormap"]
            colormap_index = 0
            for series in chart.series:
                fill = series.format.fill  # fill the legend as well
                fill.solid()
                fill.fore_color.rgb = RGBColor(colormap[colormap_index % len(colormap)][0],
                                               colormap[colormap_index % len(colormap)][1],
                                               colormap[colormap_index % len(colormap)][2])
                colormap_index += 1
        return

    def linechartformat(self, chart):
        chart.value_axis.visible = False
        chart.value_axis.has_major_gridlines = False
        chart.font.size = self.chart_format["chart_font_size"]

        # set legends
        chart.has_legend = self.chart_format["legend_bool"]
        if self.chart_format["legend_bool"]:
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.font.size = self.chart_format["legend_font_size"]
            chart.legend.include_in_layout = False

        # set data labels
        chart.plots[0].has_data_labels = self.chart_format["label_bool"]
        if self.chart_format["label_bool"]:
            labels = chart.plots[0].data_labels
            labels.number_format_is_linked = False
            labels.number_format = self.chart_format["label_number_format"]
            labels.font.size = self.chart_format["label_font_size"]

        # color
        if self.chart_format["colormap"] is None:
            color_r = 0
            color_g = 0
            color_b = 0
            colorgradient = len(chart.series)
            for series in chart.series:
                fill = series.format.fill  # fill the legend as well
                fill.solid()
                fill.fore_color.rgb = RGBColor(color_r, color_g, color_b)
                color_r += int(255/colorgradient)
                color_g += int(255/colorgradient)
                color_b += int(255/colorgradient)
        else:
            # color map format
            # [[color_R1, color_G1, color_B1], [color_R2, color_G2, color_B2],...]
            colormap = self.chart_format["colormap"]
            colormap_index = 0
            for series in chart.series:
                fill = series.format.fill  # fill the legend as well
                fill.solid()
                fill.fore_color.rgb = RGBColor(colormap[colormap_index % len(colormap)][0],
                                               colormap[colormap_index % len(colormap)][1],
                                               colormap[colormap_index % len(colormap)][2])
                colormap_index += 1
        return

    def piechartformat(self, chart):
        chart.has_legend = self.chart_format["legend_bool"]
        chart.font.size = self.chart_format["chart_font_size"]
        # set title
        if self.chart_format["chart_title"] is not None:
            chart.has_title = True
            chart.chart_title.has_text_frame = True
            chart.chart_title.text_frame.text = self.chart_format["chart_title"]
            chart.chart_title.text_frame.paragraphs[0].font.size = Pt(13)

        if self.chart_format["legend_bool"]:
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False
            chart.legend.font.size = self.chart_format["legend_font_size"]

        chart.plots[0].has_data_labels = self.chart_format["label_bool"]
        if self.chart_format["label_bool"]:
            data_labels = chart.plots[0].data_labels
            data_labels.number_format = self.chart_format["label_number_format"]
            data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END

        if self.chart_format["colormap"] is not None:
            colormap = self.chart_format["colormap"]
            colormap_index = 0
            for point in chart.series[0].points:
                fill = point.format.fill  # fill the legend as well
                fill.solid()
                fill.fore_color.rgb = RGBColor(colormap[colormap_index % len(colormap)][0],
                                               colormap[colormap_index % len(colormap)][1],
                                               colormap[colormap_index % len(colormap)][2])
                colormap_index += 1

    def chart_format_inference(self):
        chart_format = {"legend_bool": True,
                        "label_bool": True,
                        "chart_bool": True,
                        "colormap": None,
                        "legend_font_size": Pt(9),
                        "label_font_size": Pt(9),
                        "chart_font_size": Pt(9),
                        "label_number_format": "0.0"}

        if all([all(float(number_ele).is_integer() for number_ele in x[1]) for x in self.chart_series]):
            chart_format["label_number_format"] = "0"

        return chart_format