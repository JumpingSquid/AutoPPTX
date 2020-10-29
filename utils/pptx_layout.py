"""
Pptx layout contains two modules: Designer and Manager
Designer: arrange the location of objects, and pass the design structure to Manager
Manager: receive the design structure from Designer, use the data from DataPreprocessor to create the prs accordingly

Design structure (previously the config):
Every objects on the slide, will be given an unique id. Designer will arrange the location of every objects, and
provide a json format file to the manger, specifying the location of each uid
"""


from pptx.dml.color import RGBColor
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.enum.lang import MSO_LANGUAGE_ID
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches, Pt

from data2chart import ChartCreator
from data2table import TableCreator
from data2text import TextCreator



class PrsLayoutManager:

    def __init__(self,
                 presentation,
                 object_stack,
                 layout_design_container,
                 data_container):
        self.presentation = presentation

        # store the object that exists on the slide
        self.object_stack = object_stack

        # data and location for objects
        self.data_container = data_container
        self.layout_design_container = layout_design_container

        self.prs_height = self.layout_design_container.prs_height
        self.prs_width = self.layout_design_container.prs_width

        self.table_creator = TableCreator()
        self.chart_creator = ChartCreator()
        self.text_creator = TextCreator()

    def add_chart_on_slide(self, data, slide, obj_format, position, uid):
        slide = self.chart_creator.create_chart(data=data, slide=slide, obj_format=obj_format,
                                                uid=uid, position=position)
        return slide

    def add_text_on_slide(self, data, slide, obj_format, position, uid):
        slide = self.text_creator.creat_text(data=data, slide=slide, obj_format=obj_format,
                                             position=position, uid=uid)
        return slide

    def add_table_on_slide(self, data, slide, obj_format, position, uid):
        # the position of the comment
        # location should be the arguments!
        slide = self.table_creator.create_table(data=data, slide=slide, obj_format=obj_format, position=position,
                                                uid=uid)
        return slide

    def layout_execute(self):
        page_num = len(self.layout_design_container.page_object_stack)
        for page_id in range(1, page_num+1):

            bullet_slide_layout = self.presentation.slide_layouts[6]
            slide = self.presentation.slides.add_slide(bullet_slide_layout)

            for uid in self.layout_design_container.page_object_stack[page_id]:

                object_type = uid.split("_")[0]
                data = self.data_container.get_data(uid)
                position = self.layout_design_container.get_object_position(page_id, uid)

                if object_type == 'chart':
                    obj_format = self.object_stack[uid].obj_format
                    slide = self.add_chart_on_slide(data=data, slide=slide, obj_format=obj_format, position=position,
                                                    uid=uid)
                elif object_type == 'text':
                    obj_format = self.object_stack[uid].obj_format
                    slide = self.add_text_on_slide(data=data, slide=slide, obj_format=obj_format, position=position,
                                                   uid=uid)

        return None

    # set the data container
    def _set_data_container(self, data_container):
        self.data_container = data_container


class PrsLayoutDesigner:
    """
    PrsLayoutDesigner is only responsible for arranging the location of each prs_object,
     and determine the size of each object.
    It does not consider the type or data source of each object.

    In case that no pre-defined layout is given, the module will automatically design the layout.
    The design principle will be subjective, but follows some principles:
    (1) avoid overlapping
    (2) assume all objects are in a shape of square
    (3) from top to bottom: title-text-chart-text
    (4) from left to right: text-chart

    After calculating all the location and size, the information will be stored in the object: LayoutDesignContainer.
    Then LayoutDesignContainer will be passed to PrsLayoutManager to build the final presentation file.
    """
    def __init__(self, prs_config, prs_object_pool, page_stack):
        self.data = None

        # presentation format
        self.prs_width = prs_config['prs_width']
        self.prs_height = prs_config['prs_height']

        # parameter
        self.chart_num_on_slide_constraint = 3
        self.custom_layout_flag = False

        # private variable
        self._object_pool = prs_object_pool
        self._page_stack = page_stack
        self._layout_design = LayoutDesignContainer(self.prs_width, self.prs_height, max(page_stack))
        self._uid_position_stack = {}

    def execute(self):
        for page_id in self._page_stack:
            obj_in_page = self._page_stack[page_id]
            if len(obj_in_page) > 0:
                for uid in obj_in_page:
                    obj = self._object_pool[uid]
                    if obj.position is not None:
                        position = self.position_transform(obj.position)
                    else:
                        position = (0, 0, 5, 5)
                    self._uid_position_stack[uid] = position
                    self._layout_design.add_object_on_slide(slide_page=page_id,
                                                            uid=uid,
                                                            position=position)

    def custom_layout_scan(self):
        # BETA: allow for customized layout for charts, but require the user to provide fully defined layout
        # the user should provide a list contains (1) origin and (2) the box size for each chart
        # maybe move to an independent module if the analysis is too complicated
        self.custom_layout_flag = True

    def layout_design_export(self):
        return self._layout_design

    def position_transform(self, position_tuple: tuple):
        position_rep = position_tuple[0]
        if position_rep == 'a':
            _, x, y, w, h = position_tuple
            position = (x, y, w, h)

        elif position_rep == 'rb':
            _, x, y, w, h = position_tuple
            assert (0 <= x <= 1) and (0 <= y <= 1) and (0 <= w <= 1) and (0 <= h <= 1)
            x *= self.prs_width
            w *= self.prs_width
            y *= self.prs_height
            h *= self.prs_height
            position = (x, y, w, h)

        elif position_rep in ['rr', "rl", "ru", "rd"]:
            _, ref_uid, x, y, w, h = position_tuple
            if ref_uid in self._uid_position_stack:
                ref_x, ref_y, ref_w, ref_h = self._uid_position_stack[ref_uid]
            else:
                raise ValueError(f"reference object {ref_uid} is not found")

            # four relative position representation to objects
            # rr for right bound, rl for left, ru for upper, rd for bottom
            if position_rep == 'rr':
                position = (min(x + ref_x + ref_w, self.prs_width), ref_y, w, h)
            elif position_rep == 'rl':
                position = (max(ref_x - x, 0), ref_y, w, h)
            elif position_rep == 'ru':
                position = (ref_x, max(ref_y - y, 0), w, h)
            elif position_rep == 'rd':
                position = (ref_x, min(y + ref_y + ref_h, self.prs_height), w, h)
            else:
                raise ValueError(f"unknown position representation {position_rep}")

        else:
            raise ValueError(f"unknown position representation {position_rep}")

        return position


class LayoutDesignContainer:
    def __init__(self, prs_width, prs_height, page_num):
        self.prs_width = prs_width
        self.prs_height = prs_height
        self._page_num = page_num

        self.page_object_stack = {page_id: {} for page_id in range(1, page_num+1)}

    def _set_presentation_size(self, w, h):
        self.prs_width = w
        self.prs_height = h

    def _set_page_num(self, p):
        self._page_num = p

    def add_slide_on_prs(self, slide_page: int):
        # add an empty page in the object stack, this will let the manager create an empty slide
        if slide_page in self.page_object_stack:
            print(f'page {slide_page} already exists')
        else:
            for page_id in range(slide_page):
                if page_id not in self.page_object_stack:
                    self.page_object_stack[page_id] = {}

    def add_object_on_slide(self, slide_page, uid, position):
        if slide_page not in self.page_object_stack:
            self.add_slide_on_prs(slide_page)
        self.page_object_stack[slide_page][uid] = position

    def get_object_position(self, page_id, uid):
        return self.page_object_stack[page_id][uid]
