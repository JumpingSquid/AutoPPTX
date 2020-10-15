import pandas as pd
from pptx import Presentation
from pptx.util import Inches

from utils import pptx_layout, data2comment, data2table
from preprocessing import chart_data_preprocessor


class PptxConstructor:

    def __init__(self):
        self.config = {"dataframe": None,
                       "slide": []}
        self.presentation = Presentation()
        self.presentation.slide_width = Inches(13.33)
        self.presentation.slide_height = Inches(7.5)
        self.data = None
        self.page_record = []

    def _data_read(self):
        if self.config["dataframe"][-3:] == "csv":
            self.data = pd.read_csv(self.config["dataframe"])
            print(self.config["dataframe"], "is read")
        elif self.config["dataframe"][-4:] == "xlsx":
            self.data = pd.read_excel(self.config["dataframe"])
            print(self.config["dataframe"], "is read")
        else:
            raise ValueError("Only support csv or xlsx file, please check data input format")

    def _set_presentation_size(self, width, height):
        self.presentation.slide_width = width
        self.presentation.slide_height = height

    def _set_config_datapath(self, path):
        self.config["dataframe"] = path

    def set_datapath(self, data):
        self.data = data

    def presentation_reset(self):
        self.presentation = Presentation()
        self.presentation.slide_width = Inches(13.33)
        self.presentation.slide_height = Inches(7.5)

    def slide_config_append(self, slide_config):
        if slide_config["page"] in self.page_record:
            print("WARN: Duplicate page number found, change the page number to be after the existing page")
            slide_config["page"] += 1
            new_page = slide_config["page"]
            self.page_record.append(new_page)
            self.page_record.sort()
            self.page_record = [page_id+1 for page_id in self.page_record if page_id > (new_page-1)]
        else:
            print("INFO: page", slide_config["page"], "is added")
            new_page = slide_config["page"]
            self.page_record.append(new_page)
            self.page_record.sort()
        self.config["slide"].append(slide_config)

    def slide_config_create(self, page=None):

        if page is None:
            page = max(self.page_record) + 1

        slide_config = {
                        "page": page,
                        "title": None,
                        "chart": [],
                        "table": [],
                        "comment": None,
                        "setting": {"chart_id_specifed": False}
                        }

        return slide_config

    @staticmethod
    def chart_config_create(category_col,
                            type,
                            id=None,
                            value_col=None,
                            fillna="4",
                            agg="count",
                            label_number_format="0.0",
                            rescale=[1, 1],
                            pos_indicator=None,
                            ignorena=None,
                            category_freq_rank=None,
                            colormap=None,
                            appendix=None,
                            addtotal=False,
                            charttitle=None
                            ):

        chart_config = {
            "id": id,
            "cat_columns": category_col,
            "type": type,
            "fillna": fillna,
            "val_arg": {'agg': agg,
                        "category_freq_rank": category_freq_rank,
                        "pos_indicator": pos_indicator,
                        "ignorena": ignorena,
                        "addtotal": addtotal
                        },
            "format": {"label_number_format": label_number_format,
                       "colormap": colormap,
                       "chart_title": charttitle
                       },
            "appendix": appendix,
            "rescale": rescale
        }

        if value_col is not None:
            chart_config["val_columns"] = value_col

        return chart_config

    @staticmethod
    def comment_config_create(comment_type, cat_col=None, col_alias=None, pos_indicator=None, val_col=None):
        comment_config = {"cat_col": cat_col,
                          "col_alias": col_alias,
                          "comment_type": comment_type,
                          "pos_indicator": pos_indicator,
                          "val_col": val_col}
        return comment_config

    @staticmethod
    def table_config_create():
        comment_config = {}
        return comment_config

    def pptx_save(self, path='C://Users/user/Desktop/untitled.pptx'):
        self.presentation.save(path)
        print("INFO: presentation file is saved successfully as", path)

    def pptx_execute(self):
        if self.data is None:
            self._data_read()

        # to achieve queue
        self.config["slide"].reverse()

        # produce slide which queued in the slide pool
        while len(self.config["slide"]) > 0:
            slide_config = self.config["slide"].pop()
            print("INFO: slide " + str(slide_config["page"]) + " is processed")

            # call the prslayoutmanager to arrange the position of the elements in a slide
            prslayoutmanager = pptx_layout.PrsLayoutManager(self.presentation, slide_config)

            # create comment if needed
            if slide_config["comment"] is not None:
                comment_creator = data2comment.CommentCreator(self.data, slide_config["comment"])
                self.presentation = prslayoutmanager.add_comment_on_slide(comment_creator)

            # create title if needed
            if slide_config["title"] is not None:
                self.presentation = prslayoutmanager.add_title_on_slide()

            # create chart
            if len(slide_config["chart"]) > 0:
                for chart_created_config in slide_config["chart"]:
                    chart_data_processor = chart_data_preprocessor.ChartDataProcessor(self.data, chart_created_config)
                    prslayoutmanager.read_dataprocessor(chart_data_processor)
                    self.presentation = prslayoutmanager.add_chart_on_slide()

            # create table
            if len(slide_config["table"]) > 0:
                for table_config in slide_config["table"]:
                    tablecreator = data2table.TableCreator(table_config)
                    self.presentation = prslayoutmanager.add_table_on_slide(tablecreator)

            # draw the look of the layout produced
            prslayoutmanager.layout_describe()



