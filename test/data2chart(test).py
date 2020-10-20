import unittest
import pandas as pd

from data2chart import ChartCreator, ChartFormatSetter
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.chart import XL_LEGEND_POSITION, XL_DATA_LABEL_POSITION
from pptx.util import Pt
from pptx import Presentation
from pptx.util import Inches


class TestChart(unittest.TestCase):
    def setUp(self) -> None:
        self.chart_creator = ChartCreator()


class Testinit(TestChart):
    def test_prs_initializaion(self):
        self.assertIsNotNone(self.chart_creator.prs)

    def test_format_initialization(self):
        self.assertIsNotNone(self.chart_creator.chart_format_setter)

    def test_one_slide_pool(self):
        self.assertEqual(len(self.chart_creator.slide_pool), 1)


class TestPptxCreate(TestChart):
    def test_add_slide(self):
        self.chart_creator.add_slide("main_page", 6)
        self.assertIn('main_page', self.chart_creator.slide_pool)

    def test_add_line_chart_with_no_location_input(self):
        self.chart_creator.add_slide("page_1", 6)
        df = pd.DataFrame([[1, 2, 3], [4, 5, 6], [7, 8, 9]],
                          index=[f'series{i}' for i in range(1, 4)],
                          columns=[f'category{i}' for i in range(1, 4)])
        self.chart_creator.add_chart(df, "page_1", "line")

    def test_add_bar_chart_with_no_location_input(self):
        self.chart_creator.add_slide("page_2", 6)
        df = pd.DataFrame([[1, 2, 3], [4, 5, 6], [7, 8, 9]],
                          index=[f'series{i}' for i in range(1, 4)],
                          columns=[f'category{i}' for i in range(1, 4)])
        self.chart_creator.add_chart(df, "page_2", "bar")

    def test_prs_save(self):
        # the save function will return 1 if save successfully
        output = self.chart_creator.save("new.pptx")
        self.assertEqual(output, 1)
